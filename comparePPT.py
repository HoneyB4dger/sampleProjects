from pptx import Presentation
import difflib
import re

def getTextFromPPT(prs):
    listOfPages = {}
    n = 0
    for slide in prs.slides:
        n += 1
        listOfLines = []
        for shape in slide.shapes:
            try:
                #for p in shape.text.split("\n").split("."):
                for p in re.split('\. |\n', shape.text):
                    listOfLines.append(p)
            except:
                pass
        listOfPages[slide.slide_id] = {"page": n,
                                       "title": listOfLines[0],
                                       "listOfLines": listOfLines[1:]}
    return listOfPages


f = open('C:\Users\matthews\Documents\A1.pptx', 'rb')
listOfPagesA = getTextFromPPT(Presentation(f))
f = open('C:\Users\matthews\Documents\B1.pptx', 'rb')
listOfPagesB = getTextFromPPT(Presentation(f))
f.close()

# '- '	line unique to sequence 1
# '+ '	line unique to sequence 2
# '  '	line common to both sequences
# '? '	line not present in either input sequence

listOfUpdates = []

for slide_id in listOfPagesA:
    s1 = listOfPagesA[slide_id]
    if slide_id in listOfPagesB:
        s2 = listOfPagesB[slide_id]
        diffInSlide = False
        lineUpdates = []
        print("*"*80)
        print(s2["title"])
        print(s2["listOfLines"])
        for line in difflib.ndiff(s1["listOfLines"], s2["listOfLines"]):

            if line[0] == "-" or line[0] == "+":
                if not diffInSlide:
                    if s1["title"] == s2["title"]:
                        title = s1["title"]
                    else:
                        title = s1["title"] + " > " + s2["title"]
                diffInSlide = True
                lineUpdates.append(line)
        if diffInSlide:
            listOfUpdates.append({"page": int(s1["page"]),
                              "title": title,
                              "lineUpdates": lineUpdates,
                              "pageAddOrRemove": "No"
                            })
    else:
        listOfUpdates.append({"page": int(s1["page"]),
                              "title": title,
                              "lineUpdates": [],
                              "pageAddOrRemove": "Removed"
                              })

for slide_id in listOfPagesB:
    if not slide_id in listOfPagesA:
        listOfUpdates.append({"page": int(listOfPagesB[slide_id]["page"]),
                              "title": listOfPagesB[slide_id]["title"],
                              "lineUpdates": [],
                              "pageAddOrRemove": "Added"
                              })

for u in sorted(listOfUpdates, key=lambda k: k['page']):
    print u